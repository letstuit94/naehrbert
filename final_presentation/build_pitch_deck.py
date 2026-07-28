#!/usr/bin/env python3
"""NutriWise 5-minute PITCH deck → .pptx (sage-green CI, python-pptx).

Mirrors NutriWise_Pitch_5min.html: 6 main slides + 5 appendix.
Emojis are replaced by image placeholders (dashed boxes); on the Target-Group
slide the emojis stay and persona 2 gets a photo placeholder. Vision/Mission
sit stacked to the right of Solution.

Run:  .venv/bin/python final_presentation/build_pitch_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette (matches frontend --sage tokens + persona accents)
SAGE="7C9A6A"; SAGE_DK="5F7B4F"; SAGE_LT="B7C9A9"; SAGE_SOFT="EEF2EA"
CLAY="B36A4A"; CLAY_SOFT="F3E7E0"; SAND="CFC9BA"
VIOLET="6B62C4"; VIOLET_SOFT="ECEAF7"
INK="1D1D21"; INK_SOFT="6E6F74"; CANVAS="FAFAFA"; SURFACE="FFFFFF"; LINE="E5E4E0"; WHITE="FFFFFF"
FONT="Arial"
def C(h): return RGBColor.from_string(h)
OUTDIR=os.path.dirname(os.path.abspath(__file__))

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def bg(s,h): s.background.fill.solid(); s.background.fill.fore_color.rgb=C(h)

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False,adj=None,dash=False):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
        if dash:
            ln=sp.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    sp.shadow.inherit=False
    if adj is not None:
        try: sp.adjustments[0]=adj
        except Exception: pass
    if shadow:
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'80000','dist':'30000','dir':'5400000','rotWithShape':'0'})
        clr=el.makeelement(qn('a:srgbClr'),{'val':'1D1D21'}); alp=el.makeelement(qn('a:alpha'),{'val':'12000'})
        clr.append(alp); sh.append(clr); ef.append(sh); el.append(ef)
    return sp

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=3,ls=1.0,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,col,b,it) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
            r.font.color.rgb=C(col); r.font.name=FONT
    return tb
def R(t,sz,col=INK,b=False,it=False): return (t,sz,col,b,it)
def notes(s,who,body): s.notes_slide.notes_text_frame.text=f"[{who}]\n{body}"

def kicker(s,x,y,t): txt(s,x,y,11,0.3,[[R(t.upper(),12.5,SAGE_DK,True,False)]])
def heading(s,x,y,w,t,size=30): txt(s,x,y,w,1.0,[[R(t,size,INK,True,False)]],ls=1.0)
def rule(s,x,y,w=1.55): rect(s,x,y,w,0.05,fill=SAGE)
def logo(s,x,y):
    rect(s,x,y,0.34,0.34,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,x+0.46,y-0.03,4,0.42,[[R("NUTRIWISE",13,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
def footer(s,page):
    rect(s,0.6,7.05,0.16,0.16,fill=SAGE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,0.82,7.02,4,0.3,[[R("NutriWise",9,INK_SOFT,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if page: txt(s,10.8,7.02,1.9,0.3,[[R(page,9,INK_SOFT,False,page=='appendix')]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
def bullets(s,x,y,w,items,size=13,gap=0.36,mark="›",mcol=SAGE,col=INK):
    for i,it in enumerate(items):
        yy=y+i*gap
        txt(s,x,yy-0.02,0.25,0.3,[[R(mark,size,mcol,True,False)]])
        txt(s,x+0.28,yy,w-0.28,0.5,[[R(it,size,col,False,False)]],ls=1.02)
def imgph(s,x,y,w,h,label="Bild",lab_sz=8.5):
    rect(s,x,y,w,h,fill=SAGE_SOFT,line=SAGE_LT,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.14,dash=True)
    txt(s,x,y,w,h,[[R(label,lab_sz,SAGE_DK,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def card(s,x,y,w,h,fill=SURFACE,line=LINE):
    return rect(s,x,y,w,h,fill=fill,line=line,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=True,adj=0.045)

# ================= S1 Title
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE)
rect(s,8.9,0.14,4.433,7.36,fill=SAGE_SOFT)
logo(s,0.9,0.85)
txt(s,0.9,2.05,7.4,2.4,
    [[R("Upload a grocery receipt.",40,INK,True,False)],
     [R("Get one smarter thing",40,SAGE_DK,True,False)],
     [R("to buy next time.",40,SAGE_DK,True,False)]],ls=1.05)
rule(s,0.92,4.5,3.4)
txt(s,0.92,4.75,7,0.5,[[R("A receipt-first nutrition assistant.",15,INK_SOFT,False,True)]])
txt(s,0.9,5.7,7.6,0.4,[[R("AI Product Management Capstone · Pitch · 5 minutes",12.5,INK_SOFT,False,False)]])
# team avatars
for i,(ini,nm,rl,col) in enumerate([("JR","Jennifer Rake","Product & business",SAGE),("SK","Stuart Kasemeier","Data & model",SAGE_DK)]):
    yy=2.35+i*1.45
    rect(s,9.4,yy,1.0,1.0,fill=col,shape=MSO_SHAPE.OVAL)
    txt(s,9.4,yy,1.0,1.0,[[R(ini,26,WHITE,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,10.55,yy+0.14,2.6,0.8,[[R(nm,15,INK,True,False)],[R(rl,12,INK_SOFT,False,False)]],space_after=2,anchor=MSO_ANCHOR.MIDDLE)
txt(s,9.4,5.55,3.4,0.3,[[R("Placeholder avatars — send photos to drop in.",9.5,INK_SOFT,False,True)]])
notes(s,"Both — Jennifer opens","Open on the slogan. Introduce the two of you. 5-minute version; personas & business model in the appendix.")

# ================= S2 Problem + Solution + V/M
s=slide(); bg(s,CANVAS); kicker(s,0.9,0.55,"Problem & solution"); heading(s,0.9,0.9,11,"The gap we close"); rule(s,0.92,1.62)
cx=[0.9,5.0]; cw=3.85; cy=2.05; ch=3.9
# problem
card(s,cx[0],cy,cw,ch)
imgph(s,cx[0]+0.35,cy+0.3,0.6,0.6,"Bild",8)
txt(s,cx[0]+0.35,cy+1.05,cw-0.6,0.3,[[R("PROBLEM",12.5,SAGE_DK,True,False)]])
txt(s,cx[0]+0.35,cy+1.35,cw-0.6,0.5,[[R("Logging kills the habit",21,INK,True,False)]])
bullets(s,cx[0]+0.35,cy+2.15,cw-0.6,["Every app = type in every meal","Habit dies in weeks → people quit","No effortless record of eating"])
# solution
card(s,cx[1],cy,cw,ch,fill=SAGE_SOFT,line=SAGE_SOFT)
imgph(s,cx[1]+0.35,cy+0.3,0.6,0.6,"Bild",8)
txt(s,cx[1]+0.35,cy+1.05,cw-0.6,0.3,[[R("SOLUTION",12.5,SAGE_DK,True,False)]])
txt(s,cx[1]+0.35,cy+1.35,cw-0.6,0.5,[[R("The receipt logs for you",21,INK,True,False)]])
bullets(s,cx[1]+0.35,cy+2.1,cw-0.6,["Snap receipt → pantry auto-fills","Tap what you ate — 1 tap","Weekly gaps + 1 smart pick","Honest estimates, not medical advice"],gap=0.44)
# V/M stacked right
vx=9.1; vw=3.35; vh=1.86
for i,(k,t,col) in enumerate([("VISION","Eating better, zero extra effort",SAGE),("MISSION","Turn the receipt into one kind nudge",SAGE_DK)]):
    vy=cy+i*(vh+0.18)
    card(s,vx,vy,vw,vh)
    rect(s,vx,vy+0.2,0.08,vh-0.4,fill=col)
    txt(s,vx+0.32,vy+0.28,vw-0.5,0.3,[[R(k,12,col,True,False)]])
    txt(s,vx+0.32,vy+0.66,vw-0.55,1.0,[[R(t,17,INK,True,False)]],ls=1.05)
footer(s,"02 / 06")
notes(s,"Jennifer","Problem: manual logging → people quit. Solution: the receipt logs for you — snap, tap, get gaps + one pick. Vision & mission stacked on the right.")

# ================= S3 Target group
s=slide(); bg(s,CANVAS); kicker(s,0.9,0.55,"Target group"); heading(s,0.9,0.9,11,"Who it's for"); rule(s,0.92,1.62)
segs=[("🌱","Dietary Restrictor","Proof it works","Sara, 25 · vegan",
       ["Vegan/veggie/allergies",'"Vegan = default, not checkbox"',"Unsure on iron, omega-3, B12"],CLAY,CLAY_SOFT,False),
      ("⏱️","Busy Professional","No time, low friction","Tobias, 38 · 2 kids",
       ['"Just tell me what to do"',"Hates daily input","Pays for time, not complexity"],SAGE_DK,SAGE_SOFT,True),
      ("✨","Early Adopter",'Feel "enough"',"Lena, 29 · gym 3×/wk",
       ["Design quality = dealbreaker","Hates calorie counting","Wants control, not a 2nd job"],VIOLET,VIOLET_SOFT,False)]
tw=3.83; tx0=0.9; ty=2.0; th=4.15
for i,(emo,tag,h3,pers,items,acc,accs,photo) in enumerate(segs):
    x=tx0+i*(tw+0.15)
    card(s,x,ty,tw,th)
    rect(s,x+0.35,ty+0.32,0.62,0.62,fill=accs,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.25)
    txt(s,x+0.35,ty+0.3,0.62,0.62,[[R(emo,22,INK,False,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if photo:   # persona 2 photo placeholder
        rect(s,x+tw-0.95,ty+0.3,0.62,0.62,fill=SURFACE,line=SAGE_LT,lw=1.25,shape=MSO_SHAPE.OVAL,dash=True)
        txt(s,x+tw-0.95,ty+0.3,0.62,0.62,[[R("Bild",8,SAGE_DK,False,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,x+0.35,ty+1.12,1.9,0.34,fill=accs,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,x+0.35,ty+1.13,1.9,0.32,[[R(tag.upper(),9.5,acc,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.35,ty+1.58,tw-0.6,0.4,[[R(h3,19,INK,True,False)]])
    txt(s,x+0.35,ty+2.0,tw-0.6,0.3,[[R(pers,12,INK_SOFT,False,False)]])
    bullets(s,x+0.35,ty+2.5,tw-0.6,items,size=13,gap=0.4,mcol=acc)
txt(s,0.9,6.42,11.5,0.4,[[R("Core: ",14,SAGE_DK,True,False),R("health-minded adults, 25–45, in Germany.",14,INK,False,False)]],align=PP_ALIGN.CENTER)
footer(s,"03 / 06")
notes(s,"Jennifer","One core group, three flavours. Restrictor wants proof; professional wants no friction; early adopter is design-led. Full personas in appendix.")

# ================= S4 Architecture
s=slide(); bg(s,CANVAS); kicker(s,0.9,0.55,"Architecture"); heading(s,0.9,0.9,11,"From a photo to one smart pick",size=28); rule(s,0.92,1.62)
nodes=[("Capture","Receipt photo or PDF.","React · PWA"),
       ("Extract","Read text on-device — offline.","FastAPI · Tesseract"),
       ("Match","Safest match first.","verified→BLS→OFF"),
       ("Advise","Weekly gaps + next pick.","nutrition engine")]
nw=2.7; nx0=0.9; ny=2.3; nh=3.2; gap=0.28
for i,(t,sub,tag) in enumerate(nodes):
    x=nx0+i*(nw+gap)
    card(s,x,ny,nw,nh)
    imgph(s,x+0.3,ny+0.3,nw-0.6,0.85,"Bild",9)
    txt(s,x+0.3,ny+1.35,nw-0.55,0.4,[[R(t,20,INK,True,False)]])
    txt(s,x+0.3,ny+1.85,nw-0.55,0.7,[[R(sub,13.5,INK_SOFT,False,False)]],ls=1.05)
    rect(s,x+0.3,ny+nh-0.6,nw-0.6,0.36,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,x+0.3,ny+nh-0.6,nw-0.6,0.36,[[R(tag,10.5,SAGE_DK,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if i<3: txt(s,x+nw-0.02,ny+1.2,0.34,0.5,[[R("→",24,SAGE,True,False)]],align=PP_ALIGN.CENTER)
rect(s,0.9,5.9,11.53,0.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
txt(s,1.2,6.02,11,0.65,[[R("Privacy by design:  ",13,SAGE_DK,True,False),
    R("on-device OCR · per-user row-level security · GDPR export & erase. Rule-based — the AI only breaks ties.",13,INK,False,False)]],ls=1.05,anchor=MSO_ANCHOR.MIDDLE)
footer(s,"04 / 06")
notes(s,"Stuart","One flow: capture → extract (on-device Tesseract) → tiered match (verified→BLS→OFF) → advise. Rule-based; AI only breaks ties.")

# ================= S5 Outlook
s=slide(); bg(s,CANVAS); kicker(s,0.9,0.55,"Outlook"); heading(s,0.9,0.9,11.5,"What runs today — and what's next"); rule(s,0.92,1.62)
cols=[("LIVE TODAY",SAGE_DK,"The loop runs",["Receipt → pantry → gaps → pick","Daily check-off · expiry pantry","Recipes & Nutri-Coach","Multi-user · GDPR export/erase"]),
      ("NEXT · PRODUCT",SAGE,"Broaden & sharpen",["Multilingual + English food DB","Timezone-aware dates · dark mode","Best-before (MHD) sorting","Shopping list from offers/flyers"]),
      ("NEXT · GROWTH",INK_SOFT,"Partnerships",["Supermarkets & discounters","Health-insurer bonuses","Corporate health","Live testing → 10-KPI funnel"])]
cw=3.63; x0=1.05; ty=2.15
for i,(tag,col,h3,items) in enumerate(cols):
    x=x0+i*(cw+0.24)
    rect(s,x,ty+0.98,cw,0.05,fill=col if col!=INK_SOFT else LINE)
    rect(s,x,ty+0.9,0.24,0.24,fill=col,shape=MSO_SHAPE.OVAL)
    txt(s,x,ty,cw,0.3,[[R(tag,12.5,SAGE_DK if col!=INK_SOFT else INK_SOFT,True,False)]])
    txt(s,x,ty+0.32,cw,0.4,[[R(h3,18,INK,True,False)]])
    bullets(s,x,ty+1.4,cw,items,size=13.5,gap=0.5,mcol=col if col!=INK_SOFT else SAGE)
rect(s,0.9,5.85,11.53,0.85,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.08)
txt(s,1.2,5.97,11,0.65,[[R("The bet:  ",13,SAGE_DK,True,False),
    R("do people trust a receipt-based tip enough to change their next cart — and does the daily tap stick?",13,INK,False,False)]],ls=1.05,anchor=MSO_ANCHOR.MIDDLE)
footer(s,"05 / 06")
notes(s,"Stuart","Live today: the loop runs. Next product: multilingual, dark mode, MHD sorting, offer-matched shopping list. Growth: retailers, insurers, employers.")

# ================= S6 Closing
s=slide(); bg(s,CANVAS)
rect(s,0,0,13.333,0.14,fill=SAGE); rect(s,0,7.36,13.333,0.14,fill=SAGE)
logo(s,0.9,0.85)
txt(s,0.9,2.35,11.5,1.1,[[R("Thank you.",48,INK,True,False)]])
rule(s,0.92,3.5,3.4)
txt(s,0.92,3.75,11,0.5,[[R("Upload a grocery receipt. Get one smarter thing to buy next time.",18,SAGE_DK,False,True)]])
cts=[("TEAM","Jennifer Rake & Stuart Kasemeier"),("EMAIL","jennifer.rake@arcor.de   · Stuart: [add email]"),("PROJECT & CODE","github.com/letstuit94/naehrbert")]
for i,(k,v) in enumerate(cts):
    x=0.9+i*4.1
    txt(s,x,4.75,3.9,0.3,[[R(k,11.5,INK_SOFT,True,False)]])
    txt(s,x,5.08,3.9,0.5,[[R(v,13,INK,True,False)]],ls=1.05)
txt(s,0.9,6.1,11.5,0.3,[[R("Appendix: personas · tech stack · business model — Questions welcome",12.5,INK_SOFT,False,False)]])
notes(s,"Both","Close on the slogan, give contact + repo, invite questions. Appendix has personas, tech stack, BMC.")

# ================= A1-A3 Personas
def persona(ini,name,meta,tag,acc,accs,cells,secret):
    s=slide(); bg(s,CANVAS); kicker(s,0.9,0.5,"Appendix · Persona")
    rect(s,0.9,1.0,0.75,0.75,fill=accs,shape=MSO_SHAPE.OVAL)
    txt(s,0.9,1.0,0.75,0.75,[[R(ini,20,acc,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.8,1.02,7.5,0.8,[[R(name,25,INK,True,False)],[R(meta,13,INK_SOFT,False,False)]],space_after=2)
    rect(s,10.6,1.14,1.85,0.46,fill=accs,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.5)
    txt(s,10.6,1.15,1.85,0.44,[[R(tag,11.5,acc,True,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    gx=[0.9,6.75]; gy=[2.1,3.9]; gw=5.6
    for j,(lab,items) in enumerate(cells):
        x=gx[j%2]; y=gy[j//2]
        txt(s,x,y,gw,0.3,[[R(lab.upper(),11.5,INK_SOFT,True,False)]])
        bullets(s,x,y+0.35,gw,items,size=12.5,gap=0.34,mark="·",mcol=acc)
    rect(s,0.9,5.75,11.53,0.95,fill=accs,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.06)
    txt(s,1.2,5.85,11,0.3,[[R("SECRET MOTIVATION",10.5,acc,True,False)]])
    txt(s,1.2,6.12,11,0.5,[[R(secret,13.5,INK,False,False)]],ls=1.05)
    footer(s,"appendix")
    return s

s=persona("SB","Sara Bauer, 25","Master's student · Munich · Vegan 2 years","Dietary Restrictor",CLAY,CLAY_SOFT,
  [("Daily workflow",["Cooks most meals (ethics + budget)","~8-dish rotation, narrow","Reads labels · B12 (unsure of dose)"]),
   ("Psychological drivers",["Diet = identity & values",'"Am I actually healthy?"',"Reacts badly to preachy apps"]),
   ("Decision-making",["Researches ethics + privacy policies","Active in vegan groups & Reddit",'Rejects "vegan filter" apps']),
   ("Biggest frustrations",['"Vegan = checkbox, not default"',"Constant chicken protein recs","Iron/omega-3/zinc — contradictory answers"])],
  'Shut down "but where do you get your protein?" — with data proving she\'s healthier than they are.')
notes(s,"Appendix — persona 1","Sara, the dietary restrictor. Vegan by values; wants proof she's healthy; hates preachy apps.")

s=persona("TK","Tobias Kramer, 38","Project Manager · Hamburg · Married, 2 kids, no time","Busy Professional",SAGE_DK,SAGE_SOFT,
  [("Daily workflow",["Out 7:30–19:00 · kids need dinner","Leftovers or fast food; desk sandwich","Cooks only weekends (2 of 14 meals)"]),
   ("Psychological drivers",["Good role model, no life overhaul","Runs work life in Notion — likes systems","Would engage if friction is low"]),
   ("Decision-making",["Buys on trusted recommendation","Too busy to research","Pays for time — complexity is the barrier"]),
   ("Biggest frustrations",['"No 2 hours to figure out eating well"',"Quit MyFitnessPal after 4 days","Hates apps needing daily input"])],
  "Scared of getting sick — dad's heart attack at 55. The cholesterol comment landed hard. Wants to hear he'll be fine.")
notes(s,"Appendix — persona 2","Tobias, the busy professional. No time, responds to systems. Secret: scared of getting sick like his dad.")

s=persona("LM","Lena Müller, 29","UX Designer · Berlin · Lives alone, gym 3×/week","Early Adopter",VIOLET,VIOLET_SOFT,
  [("Daily workflow",["Skips breakfast · lunch near office","Tracks workouts in Strava, never a macro","Meal-preps ~monthly · saves unused recipes"]),
   ("Psychological drivers",["Control without a second job",'"Someone who works out" = identity',"Nutrition feels overwhelming, shameful"]),
   ("Decision-making",["Tries 3 apps free before committing","Reads reviews obsessively","Ugly app = instant no"]),
   ("Biggest frustrations",['"Eat healthy but always tired"',"Hates calorie counting","Apps feel built for bodybuilders / ED"])],
  "Stop feeling guilty about food and still look good. An app that makes her feel she's doing enough is the win.")
notes(s,"Appendix — persona 3","Lena, the early adopter. Design-led; nutrition feels shameful. Secret: stop feeling guilty, feel she's doing enough.")

# ================= A4 Tech stack
s=slide(); bg(s,CANVAS); kicker(s,0.9,0.55,"Appendix · Tech stack"); heading(s,0.9,0.9,11,"What it's built with",size=28); rule(s,0.92,1.62)
rows=[("Frontend","React 19 · React Router 7 · TypeScript · Vite 8"),
      ("Backend","FastAPI · Uvicorn · Pydantic v2 · Python 3.11"),
      ("OCR","Tesseract · PyMuPDF · OpenCV · on-device"),
      ("Matching / data","BLS 4.0 · OpenFoodFacts · tiered resolver"),
      ("Storage","Supabase · PostgreSQL · Row-Level Security"),
      ("AI (optional)","Gemini 2.5 Flash · recipes + coach only"),
      ("Tooling / CI","pytest · ruff · ESLint · Prettier · pre-commit · GitHub Actions")]
ry=2.1
for i,(lab,val) in enumerate(rows):
    y=ry+i*0.66
    txt(s,0.9,y,2.9,0.5,[[R(lab,14,INK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,3.9,y+0.06,8.5,0.44,fill=SAGE_SOFT,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.3)
    txt(s,4.15,y+0.06,8.1,0.44,[[R(val,12.5,SAGE_DK,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,0.9,y+0.62,11.5,0.012,fill=LINE)
footer(s,"appendix")
notes(s,"Appendix — tech stack","React/Vite front, FastAPI back, on-device Tesseract, tiered resolver over BLS+OFF, Supabase Postgres w/ RLS, Gemini only for recipes/coach.")

# ================= A5 Business Model Canvas
s=slide(); bg(s,CANVAS); kicker(s,0.9,0.5,"Appendix · Business model"); heading(s,0.9,0.85,11,"Business Model Canvas",size=27); rule(s,0.92,1.5)
def blk(x,y,w,h,title,items,hl=False,sz=8.5):
    rect(s,x,y,w,h,fill=SAGE_SOFT if hl else SURFACE,line=None if hl else LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,adj=0.05)
    txt(s,x+0.12,y+0.08,w-0.2,0.24,[[R(title.upper(),9,SAGE_DK,True,False)]])
    for i,it in enumerate(items):
        txt(s,x+0.12,y+0.36+i*0.235,w-0.22,0.24,[[R("· ",sz,SAGE,True,False),R(it,sz,INK,False,False)]],ls=0.98)
x0=0.55; cw=2.38; g=0.1; xs=[x0+i*(cw+g) for i in range(5)]
ty=1.75; hh=3.5; half=(hh-g)/2
blk(xs[0],ty,cw,hh,"Key partners",["BLS 4.0 (food table)","Open Food Facts","Hosting (Supabase)","LLM (Google Gemini)","Health insurers (future)"])
blk(xs[1],ty,cw,half,"Key activities",["Nutrient-gap analysis","Coaching & recipes","Personalization loop","GDPR data handling"])
blk(xs[1],ty+half+g,cw,half,"Key resources",["BLS 4.0 + OFF data","Tiered matching resolver","Nutrition/gap engine","2-person team"])
blk(xs[2],ty,cw,hh,"Value propositions",["Personalized, not one-size","Instant nutrient-gap analysis","Auto shopping list to goals","Respects allergies & diets","Privacy-first, fully GDPR"],hl=True)
blk(xs[3],ty,cw,half,"Customer relationships",["AI coach (chat)","Weekly check-ins & nudges","Gamification: streaks/goals","Self-service onboarding"])
blk(xs[3],ty+half+g,cw,half,"Channels",["Web app (React PWA)","Health/fitness social","Word-of-mouth","Later: app store"])
blk(xs[4],ty,cw,hh,"Customer segments",["Health-minded 25–45","Nutrition goals","Vegan/veggie/allergy","Busy professionals"],hl=True)
by=ty+hh+g; bh=1.15
blk(x0,by,cw*3+g*2,bh,"Cost structure",["Gemini API (token) · Hosting (Supabase) · Postgres · Dev (capstone 0€) · GDPR setup · Marketing (organic)"])
blk(xs[3],by,cw*2+g,bh,"Revenue streams",["Freemium → Premium · B2B insurer licensing","Affiliate grocery commissions · Corporate wellness"])
txt(s,0.55,by+bh+0.06,11.8,0.3,[[R("Reflects the built product stack — BLS 4.0 · Open Food Facts · Supabase · Gemini · React.",10.5,INK_SOFT,False,True)]])
footer(s,"appendix")
notes(s,"Appendix — business model","Value: personalized, privacy-first nutrition from the receipt. Revenue: freemium, B2B insurer licensing, grocery affiliate, corporate wellness.")

out=os.path.join(OUTDIR,"NutriWise_Pitch_5min.pptx")
prs.save(out); print("saved",out,"slides:",len(prs.slides._sldIdLst))
